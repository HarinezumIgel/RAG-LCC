# pyright: reportUnknownParameterType=false, reportMissingParameterType=false, reportUnknownVariableType=false, reportUnknownMemberType=false
"""Tests for the VisualMarkers package."""

import importlib
import io
import os
import sys
import zipfile
from pathlib import Path

import pytest

ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
SOURCE = os.path.join(ROOT, "src")
if SOURCE not in sys.path:
    sys.path.insert(0, SOURCE)

from VisualMarkers import ChunkSnippet, VisualMarkerFactory
from VisualMarkers.DocxVisualMarker import DocxVisualMarker
from VisualMarkers.PdfVisualMarker import PdfVisualMarker
from VisualMarkers.PlainTextVisualMarker import PlainTextVisualMarker
from VisualMarkers.PptxVisualMarker import PptxVisualMarker
from VisualMarkers.VisualMarker import VisualMarker

_HEDGEHOGS_PDF = Path(ROOT) / "TestDocs" / "Hedgehogs.pdf"
_APES_DOCX = Path(ROOT) / "TestDocs" / "Apes.docx"
_LIONS_PPTX = Path(ROOT) / "TestDocs" / "Lions.pptx"
_FISH_TXT = Path(ROOT) / "TestDocs" / "Fish.txt"
_CATS_MD = Path(ROOT) / "TestDocs" / "Cats.md"


# ---------------------------------------------------------------------------
# Factory
# ---------------------------------------------------------------------------


def test_factory_returns_pdf_marker_for_pdf():
    marker = VisualMarkerFactory.for_path("foo.pdf")
    assert isinstance(marker, PdfVisualMarker)


def test_factory_returns_pdf_marker_case_insensitive():
    marker = VisualMarkerFactory.for_path(Path("FOO.PDF"))
    assert isinstance(marker, PdfVisualMarker)


def test_factory_returns_docx_marker_for_docx():
    marker = VisualMarkerFactory.for_path("foo.docx")
    assert isinstance(marker, DocxVisualMarker)


def test_factory_returns_pptx_marker_for_pptx():
    marker = VisualMarkerFactory.for_path("foo.pptx")
    assert isinstance(marker, PptxVisualMarker)


def test_factory_returns_pptx_marker_for_ppt():
    marker = VisualMarkerFactory.for_path(Path("FOO.PPT"))
    assert isinstance(marker, PptxVisualMarker)


def test_factory_returns_plaintext_marker_for_txt():
    marker = VisualMarkerFactory.for_path("foo.txt")
    assert isinstance(marker, PlainTextVisualMarker)


def test_factory_returns_plaintext_marker_for_md():
    marker = VisualMarkerFactory.for_path(Path("FOO.MD"))
    assert isinstance(marker, PlainTextVisualMarker)


@pytest.mark.parametrize("name", ["foo.html", "foo"])
def test_factory_returns_none_for_unsupported(name):
    assert VisualMarkerFactory.for_path(name) is None


# ---------------------------------------------------------------------------
# PdfVisualMarker
# ---------------------------------------------------------------------------


def _has_pdfplumber() -> bool:
    try:
        import pdfplumber  # type: ignore[import-not-found]  # noqa: F401

        return True
    except ImportError:
        return False


def _has_pypdf() -> bool:
    try:
        from pypdf import PdfReader  # type: ignore[import-not-found]  # noqa: F401

        return True
    except ImportError:
        return False


pdf_backends_required = pytest.mark.skipif(
    not (_has_pdfplumber() and _has_pypdf()),
    reason="pdfplumber and pypdf not installed",
)
hedgehogs_required = pytest.mark.skipif(
    not _HEDGEHOGS_PDF.is_file(), reason="Hedgehogs.pdf test fixture missing"
)
apes_docx_required = pytest.mark.skipif(
    not _APES_DOCX.is_file(), reason="Apes.docx test fixture missing"
)
lions_pptx_required = pytest.mark.skipif(
    not _LIONS_PPTX.is_file(), reason="Lions.pptx test fixture missing"
)


def _has_python_docx() -> bool:
    try:
        import docx  # noqa: F401

        return True
    except ImportError:
        return False


def _has_python_pptx() -> bool:
    try:
        from pptx import Presentation  # noqa: F401

        return True
    except ImportError:
        return False


docx_required = pytest.mark.skipif(
    not _has_python_docx(), reason="python-docx not installed"
)
pptx_required = pytest.mark.skipif(
    not _has_python_pptx(), reason="python-pptx not installed"
)


@pdf_backends_required
@hedgehogs_required
def test_mark_to_bytes_returns_valid_pdf_bytes():
    """Highlighting a known page snippet returns parseable PDF bytes."""
    import pdfplumber  # type: ignore[import-not-found]

    with pdfplumber.open(str(_HEDGEHOGS_PDF)) as doc:
        page_text = doc.pages[0].extract_text() or ""
    snippet_text = page_text.strip().splitlines()[0].strip()
    assert len(snippet_text) >= 5  # sanity check on the fixture

    marker = PdfVisualMarker()
    out = marker.mark_to_bytes(
        _HEDGEHOGS_PDF,
        [ChunkSnippet(text=snippet_text, page_number=1)],
    )

    assert isinstance(out, (bytes, bytearray))
    assert out[:4] == b"%PDF"
    assert len(out) > 1000  # non-trivial document
    # Highlight annotations are stored as /Subtype/Highlight in the PDF.
    assert b"/Highlight" in bytes(out)


@pdf_backends_required
@hedgehogs_required
def test_mark_to_bytes_strips_page_prefix():
    """``Page N`` prefix added by PdfPageChunker must be stripped before search."""
    import pdfplumber  # type: ignore[import-not-found]

    with pdfplumber.open(str(_HEDGEHOGS_PDF)) as doc:
        page_text = doc.pages[0].extract_text() or ""
    real_line = page_text.strip().splitlines()[0].strip()
    chunk_text = f"Page 1\n\n{real_line}"

    out = PdfVisualMarker().mark_to_bytes(
        _HEDGEHOGS_PDF,
        [ChunkSnippet(text=chunk_text, page_number=1)],
    )

    assert b"/Highlight" in bytes(out)


@pdf_backends_required
@hedgehogs_required
def test_mark_to_bytes_no_match_still_returns_valid_pdf():
    """Snippet with no match must not crash; output is still valid PDF."""
    out = PdfVisualMarker().mark_to_bytes(
        _HEDGEHOGS_PDF,
        [
            ChunkSnippet(
                text="this string does not appear anywhere zzz123", page_number=1
            )
        ],
    )
    assert out[:4] == b"%PDF"


@pdf_backends_required
@hedgehogs_required
def test_mark_to_bytes_empty_snippets_returns_valid_pdf():
    out = PdfVisualMarker().mark_to_bytes(_HEDGEHOGS_PDF, [])
    assert out[:4] == b"%PDF"


@pdf_backends_required
@hedgehogs_required
def test_mark_to_bytes_falls_back_when_page_number_missing():
    """When PageNumber is None, all pages are scanned and the snippet found."""
    import pdfplumber  # type: ignore[import-not-found]

    with pdfplumber.open(str(_HEDGEHOGS_PDF)) as doc:
        page_text = doc.pages[0].extract_text() or ""
    snippet_text = page_text.strip().splitlines()[0].strip()

    out = PdfVisualMarker().mark_to_bytes(
        _HEDGEHOGS_PDF,
        [ChunkSnippet(text=snippet_text, page_number=None)],
    )

    assert b"/Highlight" in bytes(out)


@pdf_backends_required
@hedgehogs_required
def test_mark_to_bytes_does_not_modify_source(tmp_path):
    """The source PDF on disk must be byte-identical after marking."""
    work = tmp_path / "Hedgehogs.pdf"
    work.write_bytes(_HEDGEHOGS_PDF.read_bytes())
    before = work.read_bytes()

    PdfVisualMarker().mark_to_bytes(
        work,
        [ChunkSnippet(text="Hedgehog", page_number=1)],
    )

    assert work.read_bytes() == before


# ---------------------------------------------------------------------------
# Abstract base
# ---------------------------------------------------------------------------


def test_visual_marker_base_raises_not_implemented():
    with pytest.raises(NotImplementedError):
        VisualMarker().mark_to_bytes(Path("x.pdf"), [])


def test_answer_grounder_uses_singleton_config(monkeypatch):
    import Config.Config as cfg_mod

    class FakeConfig:
        def __init__(self, *args, **kwargs):
            self._initialized = True

        def get(self, key, default=None, allow_indirect=True, *, silent=False):
            if key == "_MARKED_DOCS_GROUNDING":
                return {
                    "min_sentence_tokens": 3,
                    "min_fragment_len": 7,
                    "min_overlap_window": 4,
                }
            if key == "_MARKED_DOCS_GROUNDING.min_sentence_tokens":
                return 3
            if key == "_MARKED_DOCS_GROUNDING.min_fragment_len":
                return 7
            if key == "_MARKED_DOCS_GROUNDING.min_overlap_window":
                return 4
            return default

        def get_int(self, key, default=0, *, silent=False):
            if key == "_MARKED_DOCS_GROUNDING.min_sentence_tokens":
                return 3
            if key == "_MARKED_DOCS_GROUNDING.min_fragment_len":
                return 7
            if key == "_MARKED_DOCS_GROUNDING.min_overlap_window":
                return 4
            return int(default)

    monkeypatch.setattr(cfg_mod, "Config", FakeConfig, raising=False)

    module = importlib.import_module("VisualMarkers.AnswerGrounder")
    module = importlib.reload(module)
    module.AnswerGrounder._reset()

    grounded = module.AnswerGrounder()

    assert grounded.min_sentence_tokens == 3
    assert grounded.min_fragment_len == 7
    assert grounded.min_overlap_window == 4
    assert (
        grounded.ground_answer_cli("A short answer", ["A matching chunk"])
        == "A short answer"
    )


# ---------------------------------------------------------------------------
# ImportError path for pdfplumber / pypdf
# ---------------------------------------------------------------------------


def test_pdf_marker_raises_runtime_error_when_pdfplumber_missing(monkeypatch):
    """If ``import pdfplumber`` fails, mark_to_bytes raises a clear RuntimeError."""
    import builtins

    real_import = builtins.__import__

    def fake_import(name, globals=None, locals=None, fromlist=(), level=0):
        if name == "pdfplumber":
            raise ImportError("simulated missing pdfplumber")
        return real_import(name, globals, locals, fromlist, level)

    monkeypatch.setattr(builtins, "__import__", fake_import)

    with pytest.raises(RuntimeError, match="pdfplumber"):
        PdfVisualMarker().mark_to_bytes(Path("anything.pdf"), [])


# ---------------------------------------------------------------------------
# RAGChatImpl._mark_sources — exercised without instantiating the singleton.
# ---------------------------------------------------------------------------


class _FakeDoc:
    def __init__(self, page_content: str, metadata: dict):
        self.page_content = page_content
        self.metadata = metadata


class _FakeSession:
    def __init__(self):
        self.marked_documents: list = []


class _SilentPretty:
    def write(self, *args, **kwargs):
        pass


class _StubCfg:
    def get(self, key, default=None):
        return None

    def get_str(self, key, default="") -> str:
        return default

    def get_bool(self, key, default=False):
        return default

    def get_int(self, key, default=0):
        return default

    def indirect_get(self, key, default=None, max_depth=5):
        return (default, key)


def _call_mark_sources(chosen):
    """Invoke RAGChatImpl._mark_sources unbound, with a stub *self*."""
    from Chat.RAGChatImpl import RAGChatImpl
    import os

    stub = type("Stub", (), {})()
    stub.pretty = _SilentPretty()
    stub.cfg = _StubCfg()
    # Add required methods to stub (delegate to the real implementation)
    stub._group_chunks_by_file = lambda chosen_docs: RAGChatImpl._group_chunks_by_file(
        stub, chosen_docs
    )
    stub._resolve_mark_colors = lambda: RAGChatImpl._resolve_mark_colors(stub)
    stub._build_grounded_snippets = (
        lambda session, chunks, color: RAGChatImpl._build_grounded_snippets(
            stub, session, chunks, color
        )
    )
    stub._produce_marked_bytes = (
        lambda grouped, grounded, color: RAGChatImpl._produce_marked_bytes(
            stub, grouped, grounded, color
        )
    )
    session = _FakeSession()
    RAGChatImpl._mark_sources(stub, session, chosen)  # type: ignore[arg-type]
    return session


@pdf_backends_required
@hedgehogs_required
def test_mark_sources_groups_by_source_and_stores_bytes():
    chosen = [
        _FakeDoc("Hedgehog", {"FilePath": str(_HEDGEHOGS_PDF), "PageNumber": 1}),
        _FakeDoc("nocturnal", {"FilePath": str(_HEDGEHOGS_PDF), "PageNumber": 1}),
    ]
    session = _call_mark_sources(chosen)
    assert len(session.marked_documents) == 1
    src, data = session.marked_documents[0]
    assert src == str(_HEDGEHOGS_PDF)
    assert isinstance(data, (bytes, bytearray))
    assert data[:4] == b"%PDF"


def test_mark_sources_skips_web_results():
    chosen = [
        _FakeDoc("anything", {"FilePath": "https://example.com/x", "Source": "Web"}),
    ]
    session = _call_mark_sources(chosen)
    assert session.marked_documents == []


def test_mark_sources_skips_missing_files():
    chosen = [
        _FakeDoc("anything", {"FilePath": "Z:/does/not/exist.pdf", "PageNumber": 1}),
    ]
    session = _call_mark_sources(chosen)
    assert session.marked_documents == []


@hedgehogs_required
def test_mark_sources_skips_unsupported_file_types(tmp_path):
    html = tmp_path / "notes.html"
    html.write_text("<p>hello</p>")
    chosen = [_FakeDoc("hello", {"FilePath": str(html)})]
    session = _call_mark_sources(chosen)
    assert session.marked_documents == []


@pdf_backends_required
@hedgehogs_required
def test_mark_sources_handles_invalid_page_number():
    """Non-integer ``PageNumber`` must not crash; falls back to all-page scan."""
    chosen = [
        _FakeDoc(
            "Hedgehog",
            {"FilePath": str(_HEDGEHOGS_PDF), "PageNumber": "not-a-number"},
        ),
    ]
    session = _call_mark_sources(chosen)
    assert len(session.marked_documents) == 1
    assert session.marked_documents[0][1][:4] == b"%PDF"


# ---------------------------------------------------------------------------
# DocxVisualMarker
# ---------------------------------------------------------------------------


@docx_required
@apes_docx_required
def test_docx_marker_returns_valid_docx_bytes():
    """Highlighting a known paragraph returns valid DOCX bytes with w:highlight."""
    from docx import Document

    doc = Document(str(_APES_DOCX))
    snippet_text = next(p.text for p in doc.paragraphs if len(p.text.strip()) > 30)

    out = DocxVisualMarker().mark_to_bytes(
        _APES_DOCX,
        [ChunkSnippet(text=snippet_text)],
    )

    assert isinstance(out, (bytes, bytearray))
    # DOCX is a ZIP; verify it opens correctly.
    with zipfile.ZipFile(io.BytesIO(out)) as zf:
        xml = zf.read("word/document.xml")
    assert b"w:highlight" in xml


@docx_required
@apes_docx_required
def test_docx_marker_no_match_still_returns_valid_docx():
    """Snippet with no match must not crash; output is a valid DOCX."""
    out = DocxVisualMarker().mark_to_bytes(
        _APES_DOCX,
        [ChunkSnippet(text="this phrase does not appear anywhere zzz123")],
    )
    # Must still be a valid ZIP/DOCX.
    with zipfile.ZipFile(io.BytesIO(out)):
        pass


@docx_required
@apes_docx_required
def test_docx_marker_empty_snippets_returns_valid_docx():
    out = DocxVisualMarker().mark_to_bytes(_APES_DOCX, [])
    with zipfile.ZipFile(io.BytesIO(out)):
        pass


@docx_required
@apes_docx_required
def test_docx_marker_does_not_modify_source(tmp_path):
    """The source DOCX on disk must be byte-identical after marking."""
    work = tmp_path / "Apes.docx"
    work.write_bytes(_APES_DOCX.read_bytes())
    before = work.read_bytes()

    DocxVisualMarker().mark_to_bytes(
        work,
        [ChunkSnippet(text="Gorillas")],
    )

    assert work.read_bytes() == before


# ---------------------------------------------------------------------------
# PptxVisualMarker
# ---------------------------------------------------------------------------


@pptx_required
@lions_pptx_required
def test_pptx_marker_returns_valid_pptx_bytes():
    """Highlighting a known slide paragraph returns valid PPTX bytes with a:highlight."""
    from pptx import Presentation

    prs = Presentation(str(_LIONS_PPTX))
    slide = list(prs.slides)[0]
    snippet_text = next(
        "".join(r.text for r in para.runs)
        for shape in slide.shapes
        if shape.has_text_frame
        for para in shape.text_frame.paragraphs
        if len("".join(r.text for r in para.runs).strip()) > 20
    )

    out = PptxVisualMarker().mark_to_bytes(
        _LIONS_PPTX,
        [ChunkSnippet(text=snippet_text, page_number=1)],
    )

    assert isinstance(out, (bytes, bytearray))
    # PPTX is a ZIP; verify slide XML contains the highlight element.
    with zipfile.ZipFile(io.BytesIO(out)) as zf:
        slide_names = sorted(
            n
            for n in zf.namelist()
            if n.startswith("ppt/slides/slide") and n.endswith(".xml")
        )
        slide_xml = zf.read(slide_names[0])
    assert b"a:highlight" in slide_xml


@pptx_required
@lions_pptx_required
def test_pptx_marker_no_match_still_returns_valid_pptx():
    """Snippet with no match must not crash; output is a valid PPTX."""
    out = PptxVisualMarker().mark_to_bytes(
        _LIONS_PPTX,
        [ChunkSnippet(text="this phrase does not appear anywhere zzz123")],
    )
    with zipfile.ZipFile(io.BytesIO(out)):
        pass


@pptx_required
@lions_pptx_required
def test_pptx_marker_empty_snippets_returns_valid_pptx():
    out = PptxVisualMarker().mark_to_bytes(_LIONS_PPTX, [])
    with zipfile.ZipFile(io.BytesIO(out)):
        pass


@pptx_required
@lions_pptx_required
def test_pptx_marker_does_not_modify_source(tmp_path):
    """The source PPTX on disk must be byte-identical after marking."""
    work = tmp_path / "Lions.pptx"
    work.write_bytes(_LIONS_PPTX.read_bytes())
    before = work.read_bytes()

    PptxVisualMarker().mark_to_bytes(
        work,
        [ChunkSnippet(text="Lions", page_number=1)],
    )

    assert work.read_bytes() == before


# ---------------------------------------------------------------------------
# PlainTextVisualMarker
# ---------------------------------------------------------------------------


fish_txt_required = pytest.mark.skipif(
    not _FISH_TXT.is_file(), reason="Fish.txt test fixture missing"
)
cats_md_required = pytest.mark.skipif(
    not _CATS_MD.is_file(), reason="Cats.md test fixture missing"
)


@fish_txt_required
def test_txt_marker_wraps_matching_line():
    """A matching line is wrapped with <mark> HTML tags."""
    first_content_line = next(
        l.strip()
        for l in _FISH_TXT.read_text(encoding="utf-8").splitlines()
        if l.strip()
    )
    out = (
        PlainTextVisualMarker()
        .mark_to_bytes(
            _FISH_TXT,
            [ChunkSnippet(text=first_content_line)],
        )
        .decode("utf-8")
    )
    assert f"<mark>{first_content_line}</mark>" in out


@fish_txt_required
def test_txt_marker_no_match_returns_unchanged():
    original = _FISH_TXT.read_text(encoding="utf-8")
    out = (
        PlainTextVisualMarker()
        .mark_to_bytes(
            _FISH_TXT,
            [ChunkSnippet(text="this phrase does not appear anywhere zzz123")],
        )
        .decode("utf-8")
    )
    assert out == original


@fish_txt_required
def test_txt_marker_empty_snippets_returns_unchanged():
    original = _FISH_TXT.read_text(encoding="utf-8")
    out = PlainTextVisualMarker().mark_to_bytes(_FISH_TXT, []).decode("utf-8")
    assert out == original


@fish_txt_required
def test_txt_marker_does_not_modify_source(tmp_path):
    work = tmp_path / "Fish.txt"
    work.write_bytes(_FISH_TXT.read_bytes())
    before = work.read_bytes()
    PlainTextVisualMarker().mark_to_bytes(work, [ChunkSnippet(text="fish")])
    assert work.read_bytes() == before


@cats_md_required
def test_md_marker_wraps_matching_line_with_mark_tag():
    """A matching Markdown paragraph line gets <mark>…</mark> wrapping."""
    body_line = next(
        l.strip()
        for l in _CATS_MD.read_text(encoding="utf-8").splitlines()
        if l.strip() and not l.startswith("#")
    )
    out = (
        PlainTextVisualMarker()
        .mark_to_bytes(
            _CATS_MD,
            [ChunkSnippet(text=body_line)],
        )
        .decode("utf-8")
    )
    assert f"<mark>{body_line}</mark>" in out


@cats_md_required
def test_md_marker_preserves_heading_prefix():
    """Heading lines keep their ``#`` prefix outside the mark tag."""
    heading_line = next(
        l.rstrip()
        for l in _CATS_MD.read_text(encoding="utf-8").splitlines()
        if l.startswith("## ")
    )
    heading_text = heading_line.lstrip("# ").strip()
    out = (
        PlainTextVisualMarker()
        .mark_to_bytes(
            _CATS_MD,
            [ChunkSnippet(text=heading_text)],
        )
        .decode("utf-8")
    )
    # The ## prefix must remain outside the mark tag.
    assert f"## <mark>{heading_text}</mark>" in out


@cats_md_required
def test_md_marker_does_not_modify_source(tmp_path):
    work = tmp_path / "Cats.md"
    work.write_bytes(_CATS_MD.read_bytes())
    before = work.read_bytes()
    PlainTextVisualMarker().mark_to_bytes(work, [ChunkSnippet(text="cats")])
    assert work.read_bytes() == before
