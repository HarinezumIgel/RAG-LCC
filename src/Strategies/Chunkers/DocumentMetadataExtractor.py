"""Format-aware document metadata extraction (singleton).

Harvests document-level metadata (author, title, dates, …) from source files
at load time so the ingestion pipeline can attach it to every chunk, regardless
of which chunker handles the file. Supported formats: PDF (pypdf) and the modern
Office XML formats — docx (python-docx), pptx (python-pptx), xlsx (openpyxl).

Each per-format reader returns a *raw* property dict keyed by that format's
attribute names. Truly ambiguous names are pre-normalised here (e.g. openpyxl's
``creator`` — which actually means *author* — is emitted as ``author``). The
raw dict is then mapped onto canonical chunk-metadata fields via the
``_METADATA_EXTRACTION.DOC_INFO_FIELDS`` synonym table, so formats that name the
same concept differently (``creation_date`` vs ``created``) collapse to one
field.
"""

from __future__ import annotations

import os
from typing import Any

from Commons.SingletonMixin import SingletonMixin
from Config.Config import Config
from Gui.PrettyWriter import PrettyWriter

# Office OPC core-property attribute names shared by python-docx and python-pptx.
_OFFICE_CORE_ATTRS: tuple[str, ...] = (
    "author",
    "title",
    "subject",
    "created",
    "modified",
    "last_modified_by",
    "keywords",
    "category",
    "comments",
)

# openpyxl DocumentProperties attribute → normalised raw key.
_XLSX_ATTR_MAP: dict[str, str] = {
    "creator": "author",  # openpyxl "creator" == author
    "title": "title",
    "subject": "subject",
    "created": "created",
    "modified": "modified",
    "lastModifiedBy": "last_modified_by",
    "keywords": "keywords",
    "category": "category",
}

# pypdf DocumentInformation attributes read directly (raw /Name keys are merged
# afterwards to pick up extras such as Keywords).
_PDF_ATTRS: tuple[str, ...] = (
    "author",
    "title",
    "subject",
    "creator",
    "producer",
    "creation_date",
    "modification_date",
)

# File types with a dedicated document-info reader; all others fall back to
# generic filesystem metadata.
_DOC_INFO_TYPES: frozenset[str] = frozenset({"pdf", "docx", "pptx", "xlsx"})


class DocumentMetadataExtractor(SingletonMixin):
    """Thread-safe singleton that extracts document metadata across formats."""

    def __init__(
        self, *, cfg: "Config | None" = None, pretty: "PrettyWriter | None" = None
    ) -> None:
        if self._initialized:
            return
        self._initialized = True
        self._cfg: Config = cfg or Config()
        self._pretty: PrettyWriter = pretty or PrettyWriter()

    # -- Public API ---------------------------------------------------------

    def extract(self, file_path: str, file_type: str) -> dict[str, Any]:
        """Return Chroma-safe canonical metadata for *file_path*.

        File types with a format-specific reader (PDF, docx, pptx, xlsx) yield
        rich document-info fields; all other types (images, text, csv, code,
        legacy Office, …) fall back to generic filesystem fields. Empty dict
        when extraction is disabled, the file is inaccessible, or nothing is
        present.
        """
        if not self._cfg.get_bool("_METADATA_EXTRACTION.ENABLED", False, silent=True):
            return {}
        if not file_path or not os.path.isfile(file_path):
            return {}
        file_type = (file_type or "").lower()
        if file_type in _DOC_INFO_TYPES:
            raw = self._read_raw(file_type, file_path)
            result: dict[str, Any] = (
                self._normalize(raw, "DOC_INFO_FIELDS") if raw else {}
            )
        else:
            # Everything else: universal filesystem metadata.
            raw = self._read_generic(file_path)
            result = self._normalize(raw, "GENERIC_FIELDS") if raw else {}
        # Total page count — only attach when the format has pages (default 0
        # for images/text/etc. is left off entirely so it is never displayed).
        pages = self._page_count(file_type, file_path)
        if pages > 0:
            result["Pages"] = pages
        return result

    def _page_count(self, file_type: str, file_path: str) -> int:
        """Return the document's page/slide count, or 0 when not applicable."""
        try:
            if file_type == "pdf":
                from pypdf import PdfReader

                return len(PdfReader(file_path).pages)
            if file_type == "pptx":
                from pptx import Presentation

                return len(Presentation(file_path).slides)
        except Exception:
            return 0
        return 0

    # -- Format dispatch ----------------------------------------------------

    def _read_raw(self, file_type: str, file_path: str) -> dict[str, Any]:
        try:
            if file_type == "pdf":
                return self._read_pdf(file_path)
            if file_type == "docx":
                return self._read_docx(file_path)
            if file_type == "pptx":
                return self._read_pptx(file_path)
            if file_type == "xlsx":
                return self._read_xlsx(file_path)
        except Exception as exc:
            self._pretty.write(
                "W",
                "MetadataExtract",
                f"Metadata extraction failed ({exc}). File: {file_path}",
            )
        return {}

    def _read_generic(self, file_path: str) -> dict[str, Any]:
        """Universal filesystem metadata for formats without doc properties."""
        from datetime import datetime

        try:
            st = os.stat(file_path)
        except OSError:
            return {}
        return {
            "size": st.st_size,
            "modified": datetime.fromtimestamp(st.st_mtime),
        }

    # -- Per-format readers -------------------------------------------------

    def _read_pdf(self, file_path: str) -> dict[str, Any]:
        from pypdf import PdfReader

        info = PdfReader(file_path).metadata
        if info is None:
            return {}
        raw: dict[str, Any] = {}
        for attr in _PDF_ATTRS:
            # pypdf date properties can raise on malformed values — guard each.
            try:
                value = getattr(info, attr, None)
            except Exception:
                value = None
            if value not in (None, ""):
                raw[attr] = value
        # Merge remaining "/Name" entries (e.g. /Keywords) without clobbering.
        try:
            for key, value in dict(info).items():
                name = str(key).lstrip("/").strip().lower()
                if name and value not in (None, "") and name not in raw:
                    raw[name] = value
        except Exception:
            pass
        return raw

    def _read_docx(self, file_path: str) -> dict[str, Any]:
        from docx import Document as WordDoc

        return self._read_office_core(WordDoc(file_path).core_properties)

    def _read_pptx(self, file_path: str) -> dict[str, Any]:
        from pptx import Presentation

        return self._read_office_core(Presentation(file_path).core_properties)

    def _read_office_core(self, props: Any) -> dict[str, Any]:
        raw: dict[str, Any] = {}
        for attr in _OFFICE_CORE_ATTRS:
            try:
                value = getattr(props, attr, None)
            except Exception:
                value = None
            if value not in (None, ""):
                raw[attr] = value
        return raw

    def _read_xlsx(self, file_path: str) -> dict[str, Any]:
        from openpyxl import load_workbook

        wb = load_workbook(file_path, read_only=True)
        try:
            props = wb.properties
            raw: dict[str, Any] = {}
            for attr, out_key in _XLSX_ATTR_MAP.items():
                value = getattr(props, attr, None)
                if value not in (None, ""):
                    raw[out_key] = value
            return raw
        finally:
            wb.close()

    # -- Normalisation ------------------------------------------------------

    def _normalize(self, raw: dict[str, Any], fields_key: str) -> dict[str, str]:
        fields: dict[str, Any] = self._cfg.get_dict(
            f"_METADATA_EXTRACTION.{fields_key}", {}, silent=True
        )
        if not fields:
            return {}
        low: dict[str, Any] = {str(k).lower(): v for k, v in raw.items()}
        result: dict[str, str] = {}
        for canonical, synonyms in fields.items():
            syn_seq: Any = synonyms or []
            for syn in syn_seq:
                value = low.get(str(syn).lower())
                if value is None:
                    continue
                text = str(value).strip()
                if text:
                    result[str(canonical)] = text
                    break
        return result
