import os
import time

from Config.Config import Config
from Gui.Colors import ORANGE
from Helpers.FileUtils import FileUtils
from Helpers.Helpers import Helpers
from Helpers.PerfLogger import PerfLogger
from Strategies.Chunkers.PageBasedChunker import PageBasedChunker, PageData


class SlideChunker(PageBasedChunker):
    """Structure-aware chunker for presentation files (``.pptx`` / ``.ppt``).

    Re-reads the file via *python-pptx* to recover per-slide boundaries
    that are lost during the flat text extraction in the loading pipeline.

    Each slide becomes one chunk: ``Slide N: <title>\\n\\n<body>``.
    If a single slide exceeds ``MAX_CHUNK_SIZE`` words it is split with
    ``RecursiveCharacterTextSplitter`` and each sub-chunk keeps the prefix.

    For non-PPTX file types the chunker falls back to treating the entire
    content as a single block.
    """

    def __init__(
        self,
        *,
        cfg: "Config | None" = None,
        helpers: "Helpers | None" = None,
        file_utils: "FileUtils | None" = None,
        chunker_name: str | None = None,
    ) -> None:
        super().__init__(
            cfg=cfg,
            helpers=helpers,
            file_utils=file_utils,
            chunker_name=chunker_name,
        )
        self.perf_logger: PerfLogger = PerfLogger()

    # -- PageBasedChunker interface -----------------------------------------

    def _parse_pages(
        self, file_type: str, file_path: str, content: str
    ) -> list[PageData]:
        if file_type == "pptx" and file_path and os.path.isfile(file_path):
            return self._parse_pptx(file_path)

        self._pretty.write(
            "W",
            "SlideChunker",
            f"No slide structure available for .{file_type}"
            f" — falling back to flat splitting. File: {file_path}",
            color=ORANGE,
        )
        text = content.strip()
        return [(1, "", text)] if text else []

    def _format_prefix(self, num: int, title: str) -> str:
        return f"Slide {num}: {title}" if title else f"Slide {num}"

    # -- PPTX parser --------------------------------------------------------

    @staticmethod
    def _parse_pptx(file_path: str) -> list[PageData]:
        """Extract per-slide title + body from a PPTX file."""
        from pptx import Presentation  # type: ignore[import-untyped]

        prs = Presentation(file_path)
        slides: list[PageData] = []

        for idx, slide in enumerate(prs.slides, start=1):
            title: str = ""
            body_parts: list[str] = []

            if slide.shapes.title and slide.shapes.title.text:
                title = slide.shapes.title.text.strip()

            for shape in slide.shapes:
                if not hasattr(shape, "text"):
                    continue
                text = shape.text.strip()
                if not text:
                    continue
                # Skip the title shape — already captured
                if shape == slide.shapes.title:
                    continue

                # Expand text frames to preserve bullet structure
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        line = para.text.strip()
                        if line:
                            body_parts.append(line)
                else:
                    body_parts.append(text)

            body = "\n".join(body_parts)
            if title or body:
                slides.append((idx, title, body))

        return slides
