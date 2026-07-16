"""python-pptx PPTX chunk highlighter.

Strategy
--------
* ``python-pptx`` iterates every shape's text frame on the target slide(s).
* For each retrieved snippet the ``SlideChunker`` prefix
  (``"Slide N"`` / ``"Slide N: title"``) is stripped before matching.
* The same bidirectional containment test used in ``DocxVisualMarker`` is
  applied at the paragraph level within each text frame.
* Highlighting is applied by injecting an ``<a:highlight>`` element (with a
  child ``<a:srgbClr val="RRGGBB"/>``) into each matching run's ``<a:rPr>``
  via direct lxml XML manipulation, since python-pptx does not expose a
  highlight API.
* The modified presentation is serialised in-memory; the source file is
  never touched.

Coordinate notes
----------------
PPTX/DrawingML uses the ``a:`` namespace
(``http://schemas.openxmlformats.org/drawingml/2006/main``).  The highlight
construct is::

    <a:rPr …>
        <a:highlight>
            <a:srgbClr val="FFFF00"/>
        </a:highlight>
    </a:rPr>
"""

from __future__ import annotations

import io
import re
from pathlib import Path
from typing import Any, Sequence

from VisualMarkers.VisualMarker import ChunkSnippet, VisualMarker

# Strip "Slide N" or "Slide N: title" prefix added by SlideChunker.
_SLIDE_PREFIX_RX = re.compile(r"^\s*Slide\s+\d+(?::[^\n]*)?\n+", re.IGNORECASE)

_MIN_PARA_TOKENS = 4
_MIN_FRAGMENT_LEN = 12
_PUNCT_RX = re.compile(r"[^\w\s]")

# Named highlight colours → 6-char uppercase hex strings for
# <a:srgbClr val="…">.
_NAMED_COLORS: dict[str, str] = {
    "yellow": "FFFF00",
    "green": "00FF00",
    "cyan": "00FFFF",
    "magenta": "FF00FF",
    "pink": "FFB6C1",
    "red": "FF0000",
    "blue": "0000FF",
    "orange": "FFA500",
    "lime": "80FF00",
}


def _parse_color(color: str) -> str:
    """Convert a CSS color name or ``#RRGGBB`` hex string to a 6-char
    uppercase hex string for use in ``<a:srgbClr val="…">``.
    Falls back to ``"FFFF00"`` (yellow) for unknown inputs.
    """
    c = color.strip().lower()
    if c in _NAMED_COLORS:
        return _NAMED_COLORS[c]
    if c.startswith("#"):
        hex_part = c[1:].upper()
        try:
            if len(hex_part) == 6:
                int(hex_part, 16)
                return hex_part
            if len(hex_part) == 3:
                return hex_part[0] * 2 + hex_part[1] * 2 + hex_part[2] * 2
        except ValueError:
            pass
    return "FFFF00"


class PptxVisualMarker(VisualMarker):
    """Highlight retrieved chunks in a PPTX file using python-pptx."""

    def mark_to_bytes(
        self,
        source_path: Path,
        snippets: Sequence[ChunkSnippet],
        *,
        highlight_color: str = "yellow",
    ) -> bytes:
        """Return *source_path* with *snippets* highlighted, as PPTX bytes."""
        try:
            from pptx import Presentation  # type: ignore[import-not-found]
        except ImportError as exc:
            raise RuntimeError(
                "python-pptx is required for PPTX visual marking. "
                "Install with: pip install python-pptx"
            ) from exc
        try:
            from pptx.oxml.ns import qn  # type: ignore[import-not-found]
        except ImportError as exc:
            raise RuntimeError(
                "python-pptx (oxml) is required for PPTX visual marking. "
                "Install with: pip install python-pptx"
            ) from exc
        try:
            from lxml import etree  # type: ignore[import-not-found]
        except ImportError as exc:
            raise RuntimeError(
                "lxml is required for PPTX visual marking. "
                "Install with: pip install lxml"
            ) from exc

        hex_color = _parse_color(highlight_color)
        prs = Presentation(str(source_path))
        num_slides = len(prs.slides)

        for snippet in snippets:
            text = _SLIDE_PREFIX_RX.sub("", snippet.text or "").strip()
            if not text:
                continue
            snippet_tokens = _tokenize(text)
            if not snippet_tokens:
                continue

            if (
                snippet.page_number is not None
                and 1 <= snippet.page_number <= num_slides
            ):
                slide_indices: list[int] = [snippet.page_number - 1]
            else:
                slide_indices = list(range(num_slides))

            for slide_idx in slide_indices:
                slide = prs.slides[slide_idx]
                _mark_slide(slide, snippet_tokens, text, hex_color, qn, etree)

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()


# ---------------------------------------------------------------------------
# Slide-level marking
# ---------------------------------------------------------------------------


def _mark_slide(
    slide: Any,
    snippet_tokens: list[str],
    snippet_text: str,
    hex_color: str,
    qn: Any,
    etree: Any,
) -> None:
    """Highlight paragraphs in *slide* that match *snippet_tokens*."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            para_text = "".join(run.text for run in para.runs)
            if not para_text.strip():
                continue
            para_tokens = _tokenize(para_text)
            if _para_matches(para_tokens, snippet_tokens, snippet_text):
                for run in para.runs:
                    _highlight_pptx_run(run._r, hex_color, qn, etree)


def _highlight_pptx_run(r_elem: Any, hex_color: str, qn: Any, etree: Any) -> None:
    """Inject ``<a:highlight><a:srgbClr val="…"/></a:highlight>`` into *r_elem*'s
    ``<a:rPr>``, replacing any pre-existing highlight element.
    """
    rPr = r_elem.get_or_add_rPr()
    for existing in rPr.findall(qn("a:highlight")):
        rPr.remove(existing)
    hl = etree.SubElement(rPr, qn("a:highlight"))
    srgb = etree.SubElement(hl, qn("a:srgbClr"))
    srgb.set("val", hex_color)


# ---------------------------------------------------------------------------
# Matching helpers (mirrors DocxVisualMarker)
# ---------------------------------------------------------------------------


def _para_matches(
    para_tokens: list[str],
    snippet_tokens: list[str],
    snippet_text: str,
) -> bool:
    if len(para_tokens) >= _MIN_PARA_TOKENS and _contains_sequence(
        snippet_tokens, para_tokens
    ):
        return True
    for frag in _iter_fragments(snippet_text):
        frag_tokens = _tokenize(frag)
        if frag_tokens and _contains_sequence(para_tokens, frag_tokens):
            return True
    return False


def _contains_sequence(haystack: list[str], needle: list[str]) -> bool:
    n = len(needle)
    if n == 0 or len(haystack) < n:
        return False
    for i in range(len(haystack) - n + 1):
        if haystack[i : i + n] == needle:
            return True
    return False


def _tokenize(text: str) -> list[str]:
    return _PUNCT_RX.sub(" ", text.lower()).split()


def _iter_fragments(text: str):
    seen: set[str] = set()
    for line in text.splitlines():
        stripped = line.strip()
        if len(stripped) >= _MIN_FRAGMENT_LEN and stripped not in seen:
            seen.add(stripped)
            yield stripped
    for sent in re.split(r"(?<=[\.\!\?])\s+", text):
        stripped = sent.strip()
        if len(stripped) >= _MIN_FRAGMENT_LEN and stripped not in seen:
            seen.add(stripped)
            yield stripped
